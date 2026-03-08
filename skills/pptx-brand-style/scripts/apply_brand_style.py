import sys
from pptx import Presentation
from pptx.util import Inches,Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
NAVY=RGBColor(0x23,0x16,0x8B)
TEAL=RGBColor(0x5D,0xDE,0xD2)
WHITE=RGBColor(0xFF,0xFF,0xFF)
GREY=RGBColor(0x33,0x33,0x33)
FONT="Archivo"
LOGO="/home/duane/.openclaw/workspace/skills/pptx-brand-style/assets/diversys_logo.png"
FOOTER_H=Inches(0.4)
LOGO_SZ=Inches(1.3)
def set_bg(slide,rgb):
	f=slide.background.fill
	f.solid()
	f.fore_color.rgb=rgb
def remove_old(slide):
	bad=[s for s in slide.shapes if s.shape_type in [13,17]]
	for s in bad: s._element.getparent().remove(s._element)
def add_footer(slide,sw,sh):
	s=slide.shapes.add_shape(1,0,sh-FOOTER_H,sw,FOOTER_H)
	s.fill.solid()
	s.fill.fore_color.rgb=NAVY
	s.line.fill.background()
	tf=s.text_frame
	p=tf.paragraphs[0]
	p.alignment=PP_ALIGN.CENTER
	run=p.add_run()
	run.text="Temp"
	run.font.color.rgb=WHITE
	run.font.name=FONT
	run.font.size=Pt(16)
def add_logo(slide,top=None,left=Inches(0.15)):
	if top is None: top=Inches(0.1)
	try:
		slide.shapes.add_picture(LOGO,left,top,width=LOGO_SZ,height=LOGO_SZ)
	except Exception as e:
		print(f"Logo warning: {e}")
def style_tf(tf,color,bold=None,align=None):
	for para in tf.paragraphs:
		if align is not None: para.alignment=align
		for run in para.runs:
			run.font.name=FONT
			run.font.color.rgb=color
			if bold is not None: run.font.bold=bold
def fix_positions(slide,sw,sh):
	for shape in slide.shapes:
		if not shape.is_placeholder: continue
		idx=shape.placeholder_format.idx
		if idx==0:
			shape.left=0
			shape.top=Inches(0.15)
			shape.width=sw
			shape.height=Inches(0.9)
		elif idx==1:
			shape.left=Inches(0.5)
			shape.top=Inches(1.7)
			shape.width=sw-Inches(0.8)
			shape.height=sh-Inches(2.1)-FOOTER_H
def apply_brand_layout(prs):
	sw=prs.slide_width
	sh=prs.slide_height
	for idx,slide in enumerate(prs.slides):
		layout=slide.slide_layout.name.lower()
		is_title=(idx==0) or ("title slide" in layout)
		print(f"  Slide {idx+1}: {layout}")
		remove_old(slide)
		if is_title:
			set_bg(slide,NAVY)
			for shape in slide.shapes:
				if shape.has_text_frame: style_tf(shape.text_frame,WHITE)
			add_logo(slide,top=Inches(0.15))
			continue
		set_bg(slide,WHITE)
		add_footer(slide,sw,sh)
		add_logo(slide)
		for shape in slide.shapes:
			if not shape.is_placeholder: continue
			if shape.placeholder_format.idx==0:
				style_tf(shape.text_frame,TEAL,bold=True,align=PP_ALIGN.CENTER)
			else:
				style_tf(shape.text_frame,GREY)
		fix_positions(slide,sw,sh)
if __name__=="__main__":
	path=sys.argv[1]
	prs=Presentation(path)
	apply_brand_layout(prs)
	prs.save(path)
	print(f"Saved: {path}")
