from pathlib import Path
from pptx import Presentation
from pptx.util import Inches

pptx_path = Path('/mnt/obsidian/20_PowerPoint/Presentations/AI_Support_Agent.pptx')
prs = Presentation(pptx_path)

# Add placeholder slides for diagrams (titles + notes to insert images)

for title, body in [
    ('AI Support Agent Workflow', 'Diagram file: AI_Support_Agent_Workflow.excalidraw'),
    ('Security and Trust Model', 'Diagram file: AI_Support_Agent_Security.excalidraw'),
    ('Roadmap Timeline', 'Diagram file: AI_Support_Agent_Roadmap.excalidraw'),
]:
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    tf = slide.shapes.placeholders[1].text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = body

prs.save(pptx_path)
print(str(pptx_path))
