from pathlib import Path
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Pt, Inches
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE

pptx_path = Path('/mnt/obsidian/20_PowerPoint/Presentations/AI_Support_Agent.pptx')
prs = Presentation(pptx_path)

# Brand colors
DEEP_BLUE = RGBColor(0x23, 0x16, 0x8B)
BRIGHT_BLUE = RGBColor(0x4E, 0x3D, 0xEC)
BRIGHT_SEA = RGBColor(0x5D, 0xDE, 0xD2)
SEA_GREEN = RGBColor(0x5C, 0xB6, 0xAC)
ACCENT_YELLOW = RGBColor(0xF2, 0xC9, 0x4C)
WARM_BLACK = RGBColor(0x1F, 0x1C, 0x19)
WARM_GREY = RGBColor(0x5D, 0x57, 0x4F)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

# Fonts (fallbacks if not installed)
FONT_TITLE = "Archivo"
FONT_BODY = "Fira Sans"
FONT_LABEL = "Fira Mono"

SLIDE_W = prs.slide_width
SLIDE_H = prs.slide_height


def send_to_back(shape):
    spTree = shape.element.getparent()
    spTree.remove(shape.element)
    spTree.insert(2, shape.element)


for idx, slide in enumerate(prs.slides):
    # Set white background
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = WHITE

    # Content slide defaults: curved navy panel and footer bar
    if idx != 0:
        panel_w = int(SLIDE_W * 0.48)
        panel = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
            0,
            0,
            panel_w,
            SLIDE_H,
        )
        panel.fill.solid()
        panel.fill.fore_color.rgb = DEEP_BLUE
        panel.line.color.rgb = DEEP_BLUE
        send_to_back(panel)

        footer_h = int(Inches(0.4))
        footer = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.RECTANGLE,
            0,
            SLIDE_H - footer_h,
            SLIDE_W,
            footer_h,
        )
        footer.fill.solid()
        footer.fill.fore_color.rgb = DEEP_BLUE
        footer.line.color.rgb = DEEP_BLUE
        send_to_back(footer)

    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for p in shape.text_frame.paragraphs:
            for r in p.runs:
                # Title vs body heuristic
                if shape == slide.shapes.title:
                    r.font.name = FONT_TITLE
                    r.font.size = Pt(36)
                    r.font.bold = True
                    r.font.color.rgb = DEEP_BLUE
                else:
                    r.font.name = FONT_BODY
                    r.font.size = Pt(20)
                    r.font.color.rgb = WARM_BLACK

# Add a simple accent line to title slide
slide0 = prs.slides[0]
line = slide0.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 140, SLIDE_W, Pt(2))
line.fill.solid()
line.fill.fore_color.rgb = BRIGHT_BLUE
line.line.color.rgb = BRIGHT_BLUE

prs.save(pptx_path)
print(pptx_path)
