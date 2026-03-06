from pathlib import Path
from pptx import Presentation

out_dir = Path('/mnt/obsidian/20_PowerPoint/Presentations')
out_dir.mkdir(parents=True, exist_ok=True)

prs = Presentation()

# Title slide
slide = prs.slides.add_slide(prs.slide_layouts[0])
slide.shapes.title.text = 'AI Support Agent Strategy for Diversys'
slide.placeholders[1].text = 'Date: 2026-03-05\nPresenter: Duane Leitch'

# Helper

def add_bullets(title, bullets):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    body = slide.shapes.placeholders[1].text_frame
    body.clear()
    for i, b in enumerate(bullets):
        p = body.paragraphs[0] if i == 0 else body.add_paragraph()
        p.text = b
        p.level = 0

add_bullets('Executive Summary', [
    'Why now: increasing expectations for speed and availability',
    'What it enables: consistent answers, deflection, and scale',
    'Recommended approach: phased pilot with clear KPIs'
])

add_bullets('Company Context and Fit', [
    'Diversys mission requires data integrity and compliance',
    'AI positioning aligns with operational efficiency and trust',
    'Support workflows are structured and measurable'
])

add_bullets('Objectives', [
    'After hours coverage without new headcount',
    'Faster response time and higher consistency',
    'Scale knowledge reuse across channels'
])

add_bullets('Current State', [
    'One person support team',
    'Low ticket volume, but high expectation for responsiveness',
    'Channels: email, chat, Zendesk',
    'Average first response about 20 minutes'
])

add_bullets('What the AI Support Agent Is', [
    'A support copilot for Tier 1 questions',
    'Escalates to human when confidence is low',
    'Logs every interaction for audit and review'
])

add_bullets('How It Works', [
    'Zendesk and chat intake routed to AI agent',
    'Confidence thresholds determine auto response or escalation',
    'Jira used for complex escalations and tracking'
])

add_bullets('Benefits', [
    '24x7 coverage and faster response',
    'Growth readiness without immediate hiring',
    'Consistent answers and knowledge reuse'
])

add_bullets('Security and Trust', [
    'Data minimization and approved sources only',
    'Auditability and human-in-the-loop review',
    'Role-based access and logging'
])

add_bullets('Cost Ranges and Economics', [
    'Platform and tooling costs',
    'LLM usage costs tied to volume',
    'Implementation and support versus hiring tradeoff'
])

add_bullets('Roadmap and Timeline', [
    'Phase 1: Discovery and requirements',
    'Phase 2: Pilot with limited scope',
    'Phase 3: Scale across channels'
])

add_bullets('KPIs and Governance', [
    'First response time',
    'Deflection rate',
    'CSAT',
    'Escalation accuracy'
])

add_bullets('Risks and Mitigations', [
    'Incorrect responses mitigated by confidence thresholds',
    'Security exposure mitigated by approved sources',
    'Adoption risk mitigated by phased rollout'
])

add_bullets('Decision Request', [
    'Approve pilot with defined scope',
    'Confirm go or no-go gates and success criteria'
])

add_bullets('Next Steps', [
    'Confirm KPIs and scope',
    'Select platform',
    'Start discovery'
])

# Diagram placeholder slides
for title, body in [
    ('AI Support Agent Workflow', 'Diagram file: AI_Support_Agent_Workflow.excalidraw'),
    ('Security and Trust Model', 'Diagram file: AI_Support_Agent_Security.excalidraw'),
    ('Roadmap Timeline', 'Diagram file: AI_Support_Agent_Roadmap.excalidraw'),
]:
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    tf = slide.shapes.placeholders[1].text_frame
    tf.clear()
    tf.paragraphs[0].text = body

out_path = out_dir / 'AI_Support_Agent.pptx'
prs.save(out_path)
print(out_path)
